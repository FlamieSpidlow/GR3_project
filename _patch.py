# -*- coding: utf-8 -*-
"""Đổi tiêu đề slide bìa cho khớp tên đồ án chính thức."""
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

RED     = RGBColor(0xC4, 0x12, 0x2F)
DARKRED = RGBColor(0x8C, 0x0B, 0x21)
FONT = 'Arial'

SRC = 'TheWeekend_Slide_DoAn.pptx'
OUT = SRC if not os.path.exists('~$TheWeekend_Slide_DoAn.pptx') else 'TheWeekend_Slide_DoAn_NEW.pptx'
prs = Presentation(SRC)
cover = list(prs.slides)[0]

target = None
for sh in cover.shapes:
    if sh.has_text_frame and sh.text_frame.text.strip().startswith('TheWeekend'):
        target = sh
        break
if target is None:
    raise SystemExit('Không tìm thấy textbox tiêu đề trên slide bìa')

tf = target.text_frame
tf.clear()

# Tên đồ án (ngắt dòng thủ công cho cân)
TITLE = ['Thiết kế và xây dựng ứng dụng web',
         'gợi ý địa điểm vui chơi cuối tuần',
         'cho trẻ em']
for i, line in enumerate(TITLE):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.04
    r = p.add_run(); r.text = line
    r.font.size = Pt(30); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = RED

# Dòng phụ: tên hệ thống
p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
p.space_before = Pt(14)
r = p.add_run(); r.text = 'Hệ thống TheWeekend'
r.font.size = Pt(20); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = DARKRED

prs.save(OUT)
print('SAVED', OUT)
